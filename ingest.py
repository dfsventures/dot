#!/home/joey/dot/venv/bin/python
"""
Dot Inbox Ingestion Script
--------------------------
Watches a Dropbox folder for new files, extracts content,
uses Claude to distill facts, and saves them to dot.db.

Supported: PDF, DOCX, PPTX, XLSX, CSV, TXT, MD

Run manually:     python ingest.py
Run on schedule:  add to crontab (see bottom of file)
"""

import os, io
from anthropic import Anthropic
from dotenv import load_dotenv
import dropbox as dbx_lib
from dropbox.files import WriteMode
import requests

load_dotenv(os.path.join(os.path.dirname(os.path.abspath(__file__)), '.dot.env'))

ANTHROPIC_API_KEY     = os.getenv("ANTHROPIC_API_KEY")
DROPBOX_TOKEN         = os.getenv("DROPBOX_TOKEN")
DROPBOX_REFRESH_TOKEN = os.getenv("DROPBOX_REFRESH_TOKEN")
DROPBOX_APP_KEY       = os.getenv("DROPBOX_APP_KEY")
DROPBOX_APP_SECRET    = os.getenv("DROPBOX_APP_SECRET")
TELEGRAM_TOKEN        = os.getenv("TELEGRAM_TOKEN")
TELEGRAM_CHAT_ID      = os.getenv("YOUR_TELEGRAM_USER_ID")

def notify_owner(text: str):
    if not (TELEGRAM_TOKEN and TELEGRAM_CHAT_ID):
        return
    try:
        requests.post(
            f"https://api.telegram.org/bot{TELEGRAM_TOKEN}/sendMessage",
            json={"chat_id": TELEGRAM_CHAT_ID, "text": text}, timeout=10,
        )
    except Exception as e:
        print(f"  Telegram notify failed: {e}")

client = Anthropic(api_key=ANTHROPIC_API_KEY)
dbx    = dbx_lib.Dropbox(
    oauth2_access_token=DROPBOX_TOKEN,
    oauth2_refresh_token=DROPBOX_REFRESH_TOKEN,
    app_key=DROPBOX_APP_KEY,
    app_secret=DROPBOX_APP_SECRET
)

# ── CONFIG ────────────────────────────────────────────────────────────────────
INBOX_FOLDER     = "/Dot Dump"
PROCESSED_FOLDER = "/Dot Dump/Processed"
FAILED_FOLDER    = "/Dot Dump/Failed"

# ── MEMORY (shared with agent.py) ─────────────────────────────────────────────
from memory import conn, save_memory, parse_json_array, list_deals, save_cached_doc

def _find_deal_match(fact: str, deal_names: list) -> str | None:
    """Return the first active deal company name found in the fact, or None."""
    fact_lower = fact.lower()
    for name in deal_names:
        if name.lower() in fact_lower:
            return name
    return None

def already_ingested(path: str) -> bool:
    row = conn.execute(
        "SELECT id FROM ingested_files WHERE dropbox_path = ? AND status = 'processed'",
        (path,)
    ).fetchone()
    return row is not None

def mark_ingested(filename: str, path: str, memory_count: int):
    conn.execute(
        "INSERT INTO ingested_files (filename, dropbox_path, memory_count) VALUES (?, ?, ?)",
        (filename, path, memory_count)
    )
    conn.commit()

def prior_fact_count(filename: str) -> int:
    """Facts already saved for this filename's source tag — nonzero means a crash-retry, not a fresh file."""
    tag_prefix = f"source:{filename}"
    return conn.execute(
        "SELECT COUNT(*) FROM memories WHERE tags LIKE ?", (f"{tag_prefix}%",)
    ).fetchone()[0]

# ── TEXT EXTRACTION ───────────────────────────────────────────────────────────
def extract_text(content: bytes, filename: str) -> str:
    name = filename.lower()

    if name.endswith('.pdf'):
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(io.BytesIO(content))
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n".join(pages)[:15000]
        except Exception as e:
            return f"[PDF extraction error: {e}]"

    elif name.endswith('.docx'):
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())[:15000]
        except Exception as e:
            return f"[DOCX extraction error: {e}]"

    elif name.endswith('.pptx'):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if t:
                            texts.append(t)
                    elif shape.has_table:
                        for row in shape.table.rows:
                            texts.append("\t".join(c.text.strip() for c in row.cells))
                if texts:
                    # No leading '[' — callers treat text starting with '[' as an error marker
                    slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
            return "\n\n".join(slides)[:15000]
        except Exception as e:
            return f"[PPTX extraction error: {e}]"

    elif name.endswith('.xlsx') or name.endswith('.xls'):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            sheets = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        rows.append("\t".join(str(c) if c is not None else "" for c in row))
                    if len(rows) > 200:  # cap for Claude extraction path
                        break
                if rows:
                    sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
            return "\n\n".join(sheets)[:15000]
        except Exception as e:
            return f"[XLSX extraction error: {e}]"

    elif name.endswith('.csv'):
        try:
            import csv
            text = content.decode('utf-8', errors='ignore')
            reader = csv.reader(text.splitlines())
            rows = ["\t".join(row) for row in reader if any(row)]
            return "\n".join(rows)[:15000]
        except Exception as e:
            return f"[CSV extraction error: {e}]"

    elif name.endswith(('.txt', '.md')):
        return content.decode('utf-8', errors='ignore')[:15000]

    else:
        return ""

# ── CLAUDE EXTRACTION ─────────────────────────────────────────────────────────
EXTRACTION_SYSTEM = """You extract structured facts from documents for an investor's personal AI agent.

The investor is Joey Benson-Aruna, a partner at DFS Lab (dfs.vc), focused on fintech,
digital financial services, AI, and tech across Africa.

From the document, extract every fact that would be useful to remember:
- For pitch decks: company name, founder(s), sector, geography, stage, raise amount,
  key metric(s), business model summary, notable risks or strengths, Joey's likely interest level
- For spreadsheets/CSVs: summarise what the data is tracking, extract key data points,
  notable companies, amounts, dates, patterns
- For reports/research: key findings, relevant companies or markets, notable statistics
- For meeting notes or documents: people involved, decisions made, action items, context
- For any document: anything a sharp investor would want to recall later

Return ONLY a JSON array of strings. Each string is one self-contained fact.
Be specific — include numbers, names, and dates where present.
Aim for 5–20 facts depending on document richness. If the document has no useful content, return [].

Example output:
["Acme Pay is a Nairobi-based payments infrastructure company founded by Jane Doe, raising $1M pre-seed.",
 "Acme Pay's key metric: 40,000 monthly active users as of Q1 2026.",
 "Founders have prior experience at a major African fintech and a global exchange."]"""

# Only the native-PDF call (image-only decks) asks for a transcript — the text path already
# has the text locally. Facts are required first, delimiter second, transcript last (F-27):
# a response truncated at max_tokens then loses only the transcript tail, never the facts.
TRANSCRIPT_CONTRACT = """
After the JSON array, output a line containing exactly ===TRANSCRIPT=== and then a faithful
markdown transcription of the document: every slide/page in order, headed "## Slide N", with all
visible text, numbers, table contents, and chart labels. Describe images only when they carry
information the text does not. Do not summarise or editorialise."""

PDF_EXTRACTION_SYSTEM = EXTRACTION_SYSTEM + "\n\n" + TRANSCRIPT_CONTRACT

PDF_SIZE_CAP = 24 * 1024 * 1024  # base64 inflates ~33%, API request limit is ~32MB

def compress_pdf(content: bytes) -> bytes:
    """Downsample an oversized PDF with Ghostscript (/ebook ≈ 150dpi images).
    Returns the compressed bytes, or the original on any failure."""
    import subprocess, tempfile
    try:
        with tempfile.TemporaryDirectory() as tmp:
            src = os.path.join(tmp, "in.pdf")
            dst = os.path.join(tmp, "out.pdf")
            with open(src, "wb") as f:
                f.write(content)
            subprocess.run(
                ["gs", "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.4",
                 "-dPDFSETTINGS=/ebook", "-dNOPAUSE", "-dQUIET", "-dBATCH",
                 f"-sOutputFile={dst}", src],
                check=True, timeout=300, capture_output=True,
            )
            with open(dst, "rb") as f:
                compressed = f.read()
        if 0 < len(compressed) < len(content):
            return compressed
        return content
    except Exception as e:
        print(f"  Ghostscript compression failed: {e}")
        return content

def extract_facts_from_pdf_with_claude(content: bytes, filename: str) -> tuple[list, str]:
    """Send the PDF itself to Claude — handles image-based decks with no text layer.
    API limits: ~32MB request (base64 inflates ~33%, so cap raw at 24MB), 100 pages.
    Returns (facts, transcript) — transcript is '' if the model didn't emit one
    (truncated response, or non-image PDF path never calling this at all)."""
    if len(content) > PDF_SIZE_CAP:
        print(f"  PDF too large ({len(content) / 1e6:.0f}MB) — compressing with Ghostscript...")
        content = compress_pdf(content)
        print(f"  Compressed to {len(content) / 1e6:.1f}MB")
    if len(content) > PDF_SIZE_CAP:
        print(f"  Still too large for Claude API after compression (> 24MB cap)")
        return [], ""
    try:
        import PyPDF2
        n_pages = len(PyPDF2.PdfReader(io.BytesIO(content)).pages)
        if n_pages > 100:
            print(f"  PDF has {n_pages} pages (API limit is 100)")
            return [], ""
    except Exception:
        pass  # unreadable page count — let the API decide
    import base64
    try:
        response = client.messages.create(
            model="claude-sonnet-5",
            max_tokens=8000,
            system=PDF_EXTRACTION_SYSTEM,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "document",
                        "source": {
                            "type": "base64",
                            "media_type": "application/pdf",
                            "data": base64.standard_b64encode(content).decode(),
                        },
                    },
                    {"type": "text", "text": f"Document filename: {filename}\n\nExtract facts from this document."},
                ],
            }],
        )
        raw = response.content[0].text
        head, _, transcript = raw.partition("===TRANSCRIPT===")
        facts = parse_json_array(head)
        if not facts:
            facts = parse_json_array(raw)   # model ignored the delimiter — behave exactly as before
        facts = [f for f in facts if isinstance(f, str) and len(f) > 10]
        return facts, transcript.strip()
    except Exception as e:
        print(f"  Claude PDF extraction error: {e}")
        return [], ""

def extract_facts_with_claude(text: str, filename: str) -> list:
    if not text or text.startswith('['):
        return []
    try:
        response = client.messages.create(
            model="claude-sonnet-5",
            max_tokens=2000,
            system=EXTRACTION_SYSTEM,
            messages=[{
                "role": "user",
                "content": f"Document filename: {filename}\n\nContent:\n{text}"
            }]
        )
        facts = parse_json_array(response.content[0].text)
        return [f for f in facts if isinstance(f, str) and len(f) > 10]
    except Exception as e:
        print(f"  Claude extraction error: {e}")
        return []

# ── STRUCTURED ROW INGESTION ──────────────────────────────────────────────────
# For spreadsheets/CSVs with 50+ data rows, ingest each row as a memory
# directly — no Claude call needed. Much richer than a summary.
STRUCTURED_ROW_THRESHOLD = 50

def row_to_memory(headers: list, row: tuple, filename: str) -> str:
    """Convert a spreadsheet row to a readable memory string."""
    parts = []
    for h, v in zip(headers, row):
        if v is None or str(v).strip() == "":
            continue
        # Clean up datetime objects
        if hasattr(v, 'strftime'):
            v = v.strftime('%Y-%m-%d')
        parts.append(f"{h}: {v}")
    return " | ".join(parts) if parts else ""

def ingest_structured_xlsx(content: bytes, filename: str) -> int:
    """Ingest a large XLSX row by row. Returns number of memories saved."""
    import openpyxl
    tag = f"source:{filename}"

    prior = prior_fact_count(filename)
    if prior:
        print(f"  {prior} memories already exist for this file — skipping (crash retry)")
        return prior

    total_saved = 0

    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_iter = ws.iter_rows(values_only=True)

        # First non-empty row is headers
        headers = None
        for row in rows_iter:
            if any(v is not None for v in row):
                headers = [str(h) if h is not None else f"Col{i}" for i, h in enumerate(row)]
                break

        if not headers:
            continue

        # Count data rows to decide if this sheet is worth row-by-row ingestion
        data_rows = []
        for row in rows_iter:
            if any(v is not None for v in row):
                data_rows.append(row)

        if len(data_rows) < STRUCTURED_ROW_THRESHOLD:
            # Small sheet — convert to text and let Claude handle it
            text_rows = ["\t".join(headers)]
            for row in data_rows:
                text_rows.append("\t".join(str(v) if v is not None else "" for v in row))
            text = f"[Sheet: {sheet_name}]\n" + "\n".join(text_rows)
            facts = extract_facts_with_claude(text[:15000], filename)
            for fact in facts:
                save_memory(fact, tags=tag)
            total_saved += len(facts)
            print(f"  Sheet '{sheet_name}': {len(data_rows)} rows → {len(facts)} facts via Claude")
            continue

        # Large sheet — row by row
        print(f"  Sheet '{sheet_name}': {len(data_rows)} rows → row-by-row ingestion")
        sheet_saved = 0
        for row in data_rows:
            memory = row_to_memory(headers, row, filename)
            if memory and len(memory) > 20:
                save_memory(memory, tags=tag)
                sheet_saved += 1
        conn.commit()
        print(f"    Saved {sheet_saved} memories")
        total_saved += sheet_saved

    return total_saved

def ingest_structured_csv(content: bytes, filename: str) -> int:
    """Ingest a large CSV row by row. Returns number of memories saved."""
    import csv as csv_lib
    tag = f"source:{filename}"

    prior = prior_fact_count(filename)
    if prior:
        print(f"  {prior} memories already exist for this file — skipping (crash retry)")
        return prior

    text = content.decode('utf-8', errors='ignore')
    reader = list(csv_lib.DictReader(text.splitlines()))

    if len(reader) < STRUCTURED_ROW_THRESHOLD:
        # Small — use Claude
        raw = "\t".join(reader[0].keys()) + "\n" if reader else ""
        raw += "\n".join("\t".join(str(v) for v in r.values()) for r in reader)
        facts = extract_facts_with_claude(raw[:15000], filename)
        for fact in facts:
            save_memory(fact, tags=tag)
        print(f"  CSV: {len(reader)} rows → {len(facts)} facts via Claude")
        return len(facts)

    # Large — row by row
    print(f"  CSV: {len(reader)} rows → row-by-row ingestion")
    saved = 0
    for row in reader:
        parts = [f"{k}: {v}" for k, v in row.items() if v and str(v).strip()]
        memory = " | ".join(parts)
        if memory and len(memory) > 20:
            save_memory(memory, tags=tag)
            saved += 1
    conn.commit()
    print(f"  Saved {saved} memories")
    return saved

# ── DROPBOX HELPERS ───────────────────────────────────────────────────────────
def ensure_folder(path: str):
    try:
        dbx.files_create_folder_v2(path)
    except dbx_lib.exceptions.ApiError as e:
        err = e.error
        if err.is_path() and err.get_path().is_conflict():
            return  # already exists — fine
        raise

def move_file(from_path: str, to_folder: str, filename: str):
    ensure_folder(to_folder)
    to_path = f"{to_folder}/{filename}"
    # If a file with that name already exists in destination, add a suffix
    try:
        dbx.files_move_v2(from_path, to_path, autorename=True)
    except Exception as e:
        print(f"  Move error: {e}")

def list_inbox() -> list:
    """List all files directly in INBOX_FOLDER (not subdirectories)."""
    try:
        result = dbx.files_list_folder(INBOX_FOLDER)
        files = []
        for entry in result.entries:
            # Only files, not folders
            if isinstance(entry, dbx_lib.files.FileMetadata):
                files.append(entry)
        # Handle pagination
        while result.has_more:
            result = dbx.files_list_folder_continue(result.cursor)
            for entry in result.entries:
                if isinstance(entry, dbx_lib.files.FileMetadata):
                    files.append(entry)
        return files
    except dbx_lib.exceptions.ApiError as e:
        if "path/not_found" in str(e):
            print(f"Inbox folder '{INBOX_FOLDER}' not found. Creating it...")
            ensure_folder(INBOX_FOLDER)
            return []
        raise

# ── SUPPORTED EXTENSIONS ──────────────────────────────────────────────────────
SUPPORTED = {'.pdf', '.docx', '.pptx', '.xlsx', '.xls', '.csv', '.txt', '.md'}

# ── MAIN INGESTION LOOP ───────────────────────────────────────────────────────
def run():
    print(f"Checking {INBOX_FOLDER}...")
    active_deal_names = [
        d['company'] for d in list_deals()
        if d['stage'] not in ('passed', 'invested')
    ]
    files = list_inbox()

    if not files:
        print("No files found.")
        return

    # Filter to supported types not yet ingested
    to_process = []
    for f in files:
        ext = os.path.splitext(f.name)[1].lower()
        if ext not in SUPPORTED:
            print(f"  Skipping unsupported type: {f.name}")
            continue
        if already_ingested(f.path_lower):
            print(f"  Already ingested: {f.name}")
            continue
        to_process.append(f)

    if not to_process:
        print("Nothing new to ingest.")
        return

    print(f"Found {len(to_process)} new file(s) to process.")

    # Cap owner notifications per run so a bad batch can't spam Telegram —
    # the first few failures get their own DM, the rest get folded into one
    # summary line at the end.
    fail_notify_count  = 0
    fail_notify_cap    = 5
    suppressed_failures = 0

    def _notify_failure(msg: str):
        nonlocal fail_notify_count, suppressed_failures
        if fail_notify_count < fail_notify_cap:
            notify_owner(msg)
            fail_notify_count += 1
        else:
            suppressed_failures += 1

    for entry in to_process:
        print(f"\nProcessing: {entry.name}")
        ext = os.path.splitext(entry.name)[1].lower()
        try:
            # Download
            _, response = dbx.files_download(entry.path_display)
            content = response.content
            print(f"  Downloaded: {len(content):,} bytes")

            # Route to structured ingestion for tabular files
            if ext in ('.xlsx', '.xls'):
                fact_count = ingest_structured_xlsx(content, entry.name)
                if fact_count == 0:
                    print(f"  Structured ingestion produced 0 memories. Moving to Failed.")
                    move_file(entry.path_display, FAILED_FOLDER, entry.name)
                    _notify_failure(f"⚠️ Dot Dump: couldn't process {entry.name} — moved to Failed/")
                    continue
            elif ext == '.csv':
                fact_count = ingest_structured_csv(content, entry.name)
                if fact_count == 0:
                    print(f"  Structured ingestion produced 0 memories. Moving to Failed.")
                    move_file(entry.path_display, FAILED_FOLDER, entry.name)
                    _notify_failure(f"⚠️ Dot Dump: couldn't process {entry.name} — moved to Failed/")
                    continue
            else:
                # All other types: extract text then use Claude
                text = extract_text(content, entry.name)
                if ext == '.pdf' and (not text or text.startswith('[') or len(text.strip()) < 200):
                    # Image-based PDF (designed pitch decks) — no local text layer.
                    # Send the PDF itself to Claude, which reads pages visually.
                    print(f"  Little/no text layer — sending PDF to Claude natively")
                    facts, transcript = extract_facts_from_pdf_with_claude(content, entry.name)
                    if not facts:
                        print(f"  No facts from native PDF read. Moving to Failed.")
                        move_file(entry.path_display, FAILED_FOLDER, entry.name)
                        _notify_failure(f"⚠️ Dot Dump: couldn't process {entry.name} — moved to Failed/")
                        continue
                    if transcript:
                        save_cached_doc("dropbox", entry.id, entry.content_hash, entry.name,
                                        transcript, kind="vision")
                elif not text or text.startswith('['):
                    print(f"  Could not extract text. Moving to Failed.")
                    move_file(entry.path_display, FAILED_FOLDER, entry.name)
                    _notify_failure(f"⚠️ Dot Dump: couldn't process {entry.name} — moved to Failed/")
                    continue
                else:
                    print(f"  Extracted {len(text):,} chars of text")
                    # Cache the full parse for free re-reads later (WS-11). Keyed on the
                    # Dropbox file id + content_hash — both survive the Processed/ move below.
                    save_cached_doc("dropbox", entry.id, entry.content_hash, entry.name, text)
                    facts = extract_facts_with_claude(text, entry.name)
                    if not facts:
                        print(f"  No facts extracted from text. Moving to Failed.")
                        move_file(entry.path_display, FAILED_FOLDER, entry.name)
                        _notify_failure(f"⚠️ Dot Dump: couldn't process {entry.name} — moved to Failed/")
                        continue
                prior = prior_fact_count(entry.name)
                if prior:
                    print(f"  {prior} memories already exist for this file — skipping fact save (crash retry)")
                    fact_count = prior
                else:
                    tag = f"source:{entry.name}"
                    for fact in facts:
                        deal_match = _find_deal_match(fact, active_deal_names)
                        fact_tag = f"{tag},deal:{deal_match}" if deal_match else tag
                        save_memory(fact, tags=fact_tag)
                    fact_count = len(facts)

            print(f"  Total memories saved: {fact_count}")

            # Move to Processed folder, then mark ingested — if the move fails, the file
            # stays unmarked and is retried next run (prior_fact_count above dedupes the retry)
            move_file(entry.path_display, PROCESSED_FOLDER, entry.name)
            print(f"  Moved to {PROCESSED_FOLDER}")
            mark_ingested(entry.name, entry.path_lower, fact_count)

        except Exception as e:
            print(f"  Error processing {entry.name}: {e}")
            try:
                move_file(entry.path_display, FAILED_FOLDER, entry.name)
                _notify_failure(f"⚠️ Dot Dump: couldn't process {entry.name} — moved to Failed/")
            except Exception:
                pass

    if suppressed_failures:
        notify_owner(f"⚠️ Dot Dump: {suppressed_failures} more file(s) failed and moved to Failed/ (see logs).")

    print("\nIngestion complete.")

    # Print summary of what's in memory now
    total = conn.execute("SELECT COUNT(*) FROM memories").fetchone()[0]
    print(f"Total memories in dot.db: {total}")

if __name__ == "__main__":
    # Single-instance lock: concurrent runs (cron + manual) contend on
    # ChromaDB/SQLite and can livelock. Second instance exits immediately.
    import fcntl
    _lock = open(os.path.join(os.path.dirname(os.path.abspath(__file__)), ".ingest.lock"), "w")
    try:
        fcntl.flock(_lock, fcntl.LOCK_EX | fcntl.LOCK_NB)
    except BlockingIOError:
        print("Another ingest run is already in progress — exiting.")
        raise SystemExit(0)
    run()
